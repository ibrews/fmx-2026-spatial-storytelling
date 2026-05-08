#!/usr/bin/env python3
"""
Extract all slide content from export fmx2.pptx
Extracts: text (verbatim, with shape info), images (saved with hash filenames),
video/GIF indicators, and notes.
"""

import sys
import os
import hashlib
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pptx.oxml.ns as ns

PPTX_PATH = "/Users/alex/Desktop/FMX/export fmx2.pptx"
MEDIA_DIR = "/Users/alex/git/fmx-2026-spatial-storytelling/media/fmx"
SLIDE_W_EMU = 9144000
SLIDE_H_EMU = 5143500

os.makedirs(MEDIA_DIR, exist_ok=True)

def md5hash(data):
    return hashlib.md5(data).hexdigest()

def emu_to_pct(val, dim):
    return round(val / dim * 100, 2)

def get_pos_str(shape):
    l = emu_to_pct(shape.left or 0, SLIDE_W_EMU)
    t = emu_to_pct(shape.top or 0, SLIDE_H_EMU)
    w = emu_to_pct(shape.width or 0, SLIDE_W_EMU)
    h = emu_to_pct(shape.height or 0, SLIDE_H_EMU)
    return f"L={l}% T={t}% W={w}% H={h}%"

def extract_text_from_shape(shape):
    """Extract all text from a shape, returns list of (placeholder_type, text)"""
    results = []

    # Determine shape label
    ph_type = None
    try:
        if shape.is_placeholder:
            ph = shape.placeholder_format
            ph_type = str(ph.type) if ph.type else "placeholder"
    except:
        pass

    label = shape.name or "unnamed"
    if ph_type:
        label = f"{label} [{ph_type}]"

    try:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                results.append((label, text))
    except:
        pass

    # Tables
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            table_text = "\n".join(rows)
            if table_text.strip():
                results.append((f"{label} [TABLE]", table_text))
    except:
        pass

    return results

def get_image_ext(content_type, blob):
    """Determine file extension from content type or blob magic bytes."""
    ct_map = {
        'image/png': 'png',
        'image/jpeg': 'jpg',
        'image/jpg': 'jpg',
        'image/gif': 'gif',
        'image/bmp': 'bmp',
        'image/tiff': 'tiff',
        'image/webp': 'webp',
        'image/svg+xml': 'svg',
        'video/mp4': 'mp4',
        'video/mpeg': 'mpeg',
        'video/quicktime': 'mov',
        'video/x-ms-wmv': 'wmv',
        'video/x-msvideo': 'avi',
        'image/x-emf': 'emf',
        'image/x-wmf': 'wmf',
    }
    if content_type and content_type in ct_map:
        return ct_map[content_type]
    # Fallback: detect from magic bytes
    if blob[:4] == b'\x89PNG':
        return 'png'
    if blob[:2] == b'\xff\xd8':
        return 'jpg'
    if blob[:6] in (b'GIF87a', b'GIF89a'):
        return 'gif'
    if blob[:4] == b'\x00\x00\x00\x18' or blob[4:8] == b'ftyp':
        return 'mp4'
    if blob[:4] == b'%PDF':
        return 'pdf'
    return 'bin'

def save_image(blob, slide_num, idx, content_type=None):
    """Save image blob to media dir, return (filename, ext)."""
    h = md5hash(blob)[:8]
    ext = get_image_ext(content_type, blob)
    fname = f"s{slide_num:02d}-{idx:02d}-{h}.{ext}"
    fpath = os.path.join(MEDIA_DIR, fname)
    if not os.path.exists(fpath):
        with open(fpath, 'wb') as f:
            f.write(blob)
    return fname, ext

def is_video_content_type(ct):
    if not ct:
        return False
    return ct.startswith('video/') or 'mp4' in ct or 'mpeg' in ct or 'quicktime' in ct

def process_slide(slide, slide_num):
    """Process one slide and return structured info."""
    result = {
        'slide_num': slide_num,
        'text_boxes': [],
        'images': [],
        'videos': [],
        'notes': '',
        'directives': [],
    }

    img_idx = 0

    # Process shapes
    for shape in slide.shapes:
        pos = get_pos_str(shape)

        # --- Text extraction ---
        text_entries = extract_text_from_shape(shape)
        for (label, text) in text_entries:
            result['text_boxes'].append({
                'label': label,
                'pos': pos,
                'text': text,
            })
            # Check for directives
            for line in text.split('\n'):
                if line.strip().startswith(('NOTE:', 'SLIDE:', 'TODO:', 'FIX:')):
                    result['directives'].append(line.strip())

        # --- Image extraction ---
        # Check for picture shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img = shape.image
                blob = img.blob
                ct = img.content_type
                fname, ext = save_image(blob, slide_num, img_idx, ct)
                result['images'].append({
                    'filename': fname,
                    'pos': pos,
                    'content_type': ct,
                    'is_video': is_video_content_type(ct),
                    'size_bytes': len(blob),
                    'shape_name': shape.name,
                })
                img_idx += 1
            except Exception as e:
                result['images'].append({
                    'filename': f'ERROR: {e}',
                    'pos': pos,
                    'content_type': None,
                    'is_video': False,
                    'size_bytes': 0,
                    'shape_name': shape.name,
                })

        # Check for media/video embedded in shapes (OLE/linked)
        # Video shapes often have shape_type == MSO_SHAPE_TYPE.MEDIA (6)
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.MEDIA
            try:
                # Try to get the actual media file
                shape_xml = shape._element
                # Look for video relationship
                for rel in slide.part.rels.values():
                    if 'video' in rel.reltype.lower() or 'media' in rel.reltype.lower():
                        pass  # handled below
                result['videos'].append({
                    'shape_name': shape.name,
                    'pos': pos,
                    'note': 'MEDIA shape (video/audio)',
                })
            except Exception as e:
                result['videos'].append({
                    'shape_name': shape.name,
                    'pos': pos,
                    'note': f'MEDIA shape (error: {e})',
                })

        # --- Check for grouped shapes with images ---
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                for child in shape.shapes:
                    child_pos = get_pos_str(child)
                    # Text in group
                    for (label, text) in extract_text_from_shape(child):
                        result['text_boxes'].append({
                            'label': f"[GROUP:{shape.name}] {label}",
                            'pos': child_pos,
                            'text': text,
                        })
                    # Images in group
                    if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            img = child.image
                            blob = img.blob
                            ct = img.content_type
                            fname, ext = save_image(blob, slide_num, img_idx, ct)
                            result['images'].append({
                                'filename': fname,
                                'pos': child_pos,
                                'content_type': ct,
                                'is_video': is_video_content_type(ct),
                                'size_bytes': len(blob),
                                'shape_name': f"[GROUP:{shape.name}] {child.name}",
                            })
                            img_idx += 1
                        except Exception as e:
                            pass
            except Exception as e:
                pass

    # --- Scan slide relationships for any embedded media ---
    try:
        for rel_id, rel in slide.part.rels.items():
            rt = rel.reltype
            is_vid = 'video' in rt.lower() or 'audio' in rt.lower()
            is_img = 'image' in rt.lower()

            if is_vid or is_img:
                try:
                    part = rel.target_part
                    blob = part.blob
                    ct = part.content_type

                    if is_vid or is_video_content_type(ct):
                        # Save video/gif
                        fname, ext = save_image(blob, slide_num, img_idx, ct)
                        result['videos'].append({
                            'filename': fname,
                            'content_type': ct,
                            'size_bytes': len(blob),
                            'rel_id': rel_id,
                            'note': 'embedded via relationship',
                        })
                        img_idx += 1
                    elif is_img:
                        # Check if already saved via shape traversal
                        h = md5hash(blob)[:8]
                        ext2 = get_image_ext(ct, blob)
                        expected = f"s{slide_num:02d}-"
                        already = any(h in img['filename'] for img in result['images'] if 'filename' in img)
                        if not already:
                            fname, ext = save_image(blob, slide_num, img_idx, ct)
                            result['images'].append({
                                'filename': fname,
                                'pos': 'via relationship (no shape pos)',
                                'content_type': ct,
                                'is_video': False,
                                'size_bytes': len(blob),
                                'shape_name': f'[rel:{rel_id}]',
                            })
                            img_idx += 1
                except Exception as e:
                    # External/linked media
                    try:
                        target = rel.target_ref
                        result['videos'].append({
                            'filename': None,
                            'content_type': None,
                            'size_bytes': 0,
                            'rel_id': rel_id,
                            'note': f'EXTERNAL/LINKED media: {target}',
                        })
                    except:
                        pass
    except Exception as e:
        pass

    # --- Notes ---
    try:
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text.strip()
        if notes_text:
            result['notes'] = notes_text
    except:
        pass

    return result

def print_slide(r):
    sn = r['slide_num']
    print(f"\n{'='*70}")
    print(f"SLIDE {sn}")
    print(f"{'='*70}")

    if r['text_boxes']:
        print(f"\n--- TEXT ---")
        for tb in r['text_boxes']:
            print(f"  Shape: {tb['label']}")
            print(f"  Pos:   {tb['pos']}")
            # Print text with indentation
            for line in tb['text'].split('\n'):
                print(f"  > {line}")
            print()

    if r['images']:
        print(f"--- IMAGES ---")
        for img in r['images']:
            flag = " [VIDEO/GIF]" if img.get('is_video') else ""
            print(f"  {img['filename']}{flag}")
            print(f"    Pos: {img['pos']}")
            print(f"    Shape: {img.get('shape_name','?')} | CT: {img.get('content_type','?')} | {img.get('size_bytes',0):,} bytes")

    if r['videos']:
        print(f"\n--- VIDEOS/MEDIA (relationships) ---")
        for v in r['videos']:
            fname = v.get('filename', '[no file]')
            print(f"  {fname or '[no file]'}")
            if 'content_type' in v:
                print(f"    CT: {v.get('content_type','?')} | {v.get('size_bytes',0):,} bytes")
            print(f"    Note: {v.get('note','')}")

    if r['notes']:
        print(f"\n--- NOTES ---")
        for line in r['notes'].split('\n'):
            print(f"  {line}")

    if r['directives']:
        print(f"\n--- DIRECTIVES ---")
        for d in r['directives']:
            print(f"  !! {d}")

def main():
    print(f"Loading: {PPTX_PATH}")
    prs = Presentation(PPTX_PATH)
    total = len(prs.slides)
    print(f"Total slides: {total}")
    print(f"Slide dimensions: {prs.slide_width} x {prs.slide_height} EMU")
    print(f"Media output dir: {MEDIA_DIR}")

    all_results = []
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        r = process_slide(slide, slide_num)
        all_results.append(r)
        print_slide(r)

    # Summary
    print(f"\n{'='*70}")
    print(f"SUMMARY")
    print(f"{'='*70}")
    print(f"Total slides: {total}")
    total_images = sum(len(r['images']) for r in all_results)
    total_videos = sum(len(r['videos']) for r in all_results)
    total_text = sum(len(r['text_boxes']) for r in all_results)
    print(f"Total text boxes: {total_text}")
    print(f"Total images extracted: {total_images}")
    print(f"Total video/media refs: {total_videos}")

    print(f"\nFiles saved to: {MEDIA_DIR}")
    saved = sorted(os.listdir(MEDIA_DIR))
    for f in saved:
        print(f"  {f}")

if __name__ == '__main__':
    main()
